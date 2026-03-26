"""
Document processor — read, create, and edit Microsoft Office & PDF files.

Supports:
- Word (.docx) — read text/tables/images, create, edit paragraphs
- Excel (.xlsx) — read sheets/cells/formulas, create, edit cells
- PowerPoint (.pptx) — read slides/text/notes, create, edit slides
- PDF (.pdf) — read text, create from text/HTML

All operations work with bytes (file content) so they integrate cleanly
with any BaseFileStore implementation (local, S3, GCS).

Dependencies: pip install "tova[documents]"
  - python-docx, openpyxl, python-pptx, pypdf, reportlab
"""

from __future__ import annotations

import io
import json
import logging
from typing import Any

logger = logging.getLogger(__name__)


# ── Word (.docx) ──────────────────────────────────────────────────


def read_docx(content: bytes) -> dict:
    """Extract text, tables, and metadata from a Word document.

    Returns:
        dict with: text, paragraphs, tables, metadata, headings
    """
    try:
        from docx import Document
    except ImportError:
        return {"error": "python-docx not installed. Run: pip install 'tova[documents]'"}

    doc = Document(io.BytesIO(content))

    paragraphs = []
    headings = []
    full_text = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name if para.style else ""
        entry = {"text": text, "style": style}
        paragraphs.append(entry)
        full_text.append(text)

        if "Heading" in style:
            level = style.replace("Heading ", "").replace("Heading", "1")
            try:
                level = int(level)
            except ValueError:
                level = 1
            headings.append({"text": text, "level": level})

    # Extract tables
    tables = []
    for i, table in enumerate(doc.tables):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        if rows:
            tables.append({
                "index": i,
                "headers": rows[0] if rows else [],
                "rows": rows[1:] if len(rows) > 1 else [],
                "total_rows": len(rows),
            })

    # Metadata
    props = doc.core_properties
    metadata = {
        "title": props.title or "",
        "author": props.author or "",
        "created": str(props.created) if props.created else "",
        "modified": str(props.modified) if props.modified else "",
        "subject": props.subject or "",
        "keywords": props.keywords or "",
    }

    return {
        "format": "docx",
        "text": "\n".join(full_text),
        "paragraphs": paragraphs,
        "headings": headings,
        "tables": tables,
        "metadata": metadata,
        "paragraph_count": len(paragraphs),
        "table_count": len(tables),
    }


def create_docx(
    title: str = "",
    content: list[dict] | str = "",
    tables: list[dict] | None = None,
) -> bytes:
    """Create a Word document.

    Args:
        title: Document title (added as Heading 1)
        content: Either a plain text string, or a list of dicts:
                 [{"text": "...", "style": "Heading 1|Heading 2|Normal|..."}]
        tables: Optional list of tables:
                [{"headers": ["Col1", "Col2"], "rows": [["a", "b"], ["c", "d"]]}]

    Returns:
        Document as bytes (.docx)
    """
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError:
        raise ImportError("python-docx not installed. Run: pip install 'tova[documents]'")

    doc = Document()

    if title:
        doc.add_heading(title, level=1)

    if isinstance(content, str):
        for line in content.split("\n"):
            line = line.strip()
            if line:
                doc.add_paragraph(line)
    elif isinstance(content, list):
        for item in content:
            text = item.get("text", "")
            style = item.get("style", "Normal")
            if "Heading" in style:
                level = style.replace("Heading ", "").replace("Heading", "1")
                try:
                    level = int(level)
                except ValueError:
                    level = 1
                doc.add_heading(text, level=level)
            else:
                doc.add_paragraph(text, style=style)

    if tables:
        for tbl in tables:
            headers = tbl.get("headers", [])
            rows = tbl.get("rows", [])
            if not headers:
                continue
            table = doc.add_table(rows=1, cols=len(headers))
            table.style = "Table Grid"
            for i, header in enumerate(headers):
                table.rows[0].cells[i].text = str(header)
            for row_data in rows:
                row = table.add_row()
                for i, val in enumerate(row_data):
                    if i < len(row.cells):
                        row.cells[i].text = str(val)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def edit_docx(content: bytes, edits: list[dict]) -> bytes:
    """Edit an existing Word document.

    Args:
        content: Existing .docx file bytes
        edits: List of edits:
            [{"action": "replace", "find": "old text", "replace": "new text"}]
            [{"action": "append", "text": "New paragraph", "style": "Normal"}]
            [{"action": "insert", "after": "heading text", "text": "...", "style": "Normal"}]
            [{"action": "delete", "find": "text to remove"}]

    Returns:
        Modified document as bytes
    """
    try:
        from docx import Document
    except ImportError:
        raise ImportError("python-docx not installed. Run: pip install 'tova[documents]'")

    doc = Document(io.BytesIO(content))

    for edit in edits:
        action = edit.get("action", "")

        if action == "replace":
            find_text = edit.get("find", "")
            replace_text = edit.get("replace", "")
            for para in doc.paragraphs:
                if find_text in para.text:
                    for run in para.runs:
                        if find_text in run.text:
                            run.text = run.text.replace(find_text, replace_text)

        elif action == "append":
            text = edit.get("text", "")
            style = edit.get("style", "Normal")
            if "Heading" in style:
                level = int(style.replace("Heading ", "").replace("Heading", "1"))
                doc.add_heading(text, level=level)
            else:
                doc.add_paragraph(text, style=style)

        elif action == "delete":
            find_text = edit.get("find", "")
            for para in doc.paragraphs:
                if find_text in para.text:
                    p = para._element
                    p.getparent().remove(p)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── Excel (.xlsx) ─────────────────────────────────────────────────


def read_xlsx(content: bytes) -> dict:
    """Extract data from an Excel workbook.

    Returns:
        dict with: sheets (list of sheet data), metadata
    """
    try:
        from openpyxl import load_workbook
    except ImportError:
        return {"error": "openpyxl not installed. Run: pip install 'tova[documents]'"}

    wb = load_workbook(io.BytesIO(content), data_only=True)

    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        headers = []

        for i, row in enumerate(ws.iter_rows(values_only=True)):
            str_row = [str(cell) if cell is not None else "" for cell in row]
            if i == 0:
                headers = str_row
            rows.append(str_row)

        # Summary stats for numeric columns
        col_stats = {}
        if len(rows) > 1:
            for col_idx, header in enumerate(headers):
                values = []
                for row in rows[1:]:
                    if col_idx < len(row):
                        try:
                            values.append(float(row[col_idx]))
                        except (ValueError, TypeError):
                            pass
                if values:
                    col_stats[header] = {
                        "min": min(values),
                        "max": max(values),
                        "sum": sum(values),
                        "avg": round(sum(values) / len(values), 2),
                        "count": len(values),
                    }

        sheets.append({
            "name": sheet_name,
            "headers": headers,
            "rows": rows[1:] if len(rows) > 1 else [],  # Exclude header row
            "total_rows": ws.max_row or 0,
            "total_cols": ws.max_column or 0,
            "column_stats": col_stats,
        })

    return {
        "format": "xlsx",
        "sheets": sheets,
        "sheet_names": wb.sheetnames,
        "sheet_count": len(wb.sheetnames),
    }


def create_xlsx(
    sheets: list[dict],
) -> bytes:
    """Create an Excel workbook.

    Args:
        sheets: List of sheet definitions:
            [{"name": "Sheet1", "headers": ["A", "B"], "rows": [["1", "2"], ["3", "4"]]}]

    Returns:
        Workbook as bytes (.xlsx)
    """
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment
    except ImportError:
        raise ImportError("openpyxl not installed. Run: pip install 'tova[documents]'")

    wb = Workbook()

    for i, sheet_def in enumerate(sheets):
        if i == 0:
            ws = wb.active
            ws.title = sheet_def.get("name", "Sheet1")
        else:
            ws = wb.create_sheet(sheet_def.get("name", f"Sheet{i+1}"))

        headers = sheet_def.get("headers", [])
        rows = sheet_def.get("rows", [])

        # Write headers with styling
        header_font = Font(bold=True)
        header_fill = PatternFill(start_color="D9E1F2", end_color="D9E1F2", fill_type="solid")

        for col_idx, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col_idx, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center")

        # Write data rows
        for row_idx, row_data in enumerate(rows, 2):
            for col_idx, value in enumerate(row_data, 1):
                # Try to preserve numeric types
                try:
                    val = float(value)
                    if val == int(val):
                        val = int(val)
                    ws.cell(row=row_idx, column=col_idx, value=val)
                except (ValueError, TypeError):
                    ws.cell(row=row_idx, column=col_idx, value=str(value))

        # Auto-width columns
        for col_idx in range(1, len(headers) + 1):
            max_width = len(str(headers[col_idx - 1])) if col_idx <= len(headers) else 10
            for row in rows:
                if col_idx - 1 < len(row):
                    max_width = max(max_width, len(str(row[col_idx - 1])))
            ws.column_dimensions[ws.cell(row=1, column=col_idx).column_letter].width = min(max_width + 2, 50)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def edit_xlsx(content: bytes, edits: list[dict]) -> bytes:
    """Edit an existing Excel workbook.

    Args:
        content: Existing .xlsx file bytes
        edits: List of edits:
            [{"action": "set_cell", "sheet": "Sheet1", "cell": "A1", "value": "new value"}]
            [{"action": "add_row", "sheet": "Sheet1", "values": ["a", "b", "c"]}]
            [{"action": "delete_row", "sheet": "Sheet1", "row": 5}]
            [{"action": "add_sheet", "name": "New Sheet", "headers": ["A", "B"]}]
            [{"action": "set_formula", "sheet": "Sheet1", "cell": "C1", "formula": "=SUM(A1:B1)"}]

    Returns:
        Modified workbook as bytes
    """
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ImportError("openpyxl not installed. Run: pip install 'tova[documents]'")

    wb = load_workbook(io.BytesIO(content))

    for edit in edits:
        action = edit.get("action", "")
        sheet_name = edit.get("sheet", wb.sheetnames[0] if wb.sheetnames else "Sheet1")

        if action == "set_cell":
            ws = wb[sheet_name] if sheet_name in wb.sheetnames else wb.active
            cell_ref = edit.get("cell", "A1")
            value = edit.get("value", "")
            try:
                val = float(value)
                if val == int(val):
                    val = int(val)
                ws[cell_ref] = val
            except (ValueError, TypeError):
                ws[cell_ref] = str(value)

        elif action == "set_formula":
            ws = wb[sheet_name] if sheet_name in wb.sheetnames else wb.active
            ws[edit.get("cell", "A1")] = edit.get("formula", "")

        elif action == "add_row":
            ws = wb[sheet_name] if sheet_name in wb.sheetnames else wb.active
            values = edit.get("values", [])
            ws.append(values)

        elif action == "delete_row":
            ws = wb[sheet_name] if sheet_name in wb.sheetnames else wb.active
            row_num = edit.get("row", 1)
            ws.delete_rows(row_num)

        elif action == "add_sheet":
            name = edit.get("name", "New Sheet")
            ws = wb.create_sheet(name)
            headers = edit.get("headers", [])
            for col_idx, h in enumerate(headers, 1):
                ws.cell(row=1, column=col_idx, value=h)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ── PowerPoint (.pptx) ────────────────────────────────────────────


def read_pptx(content: bytes) -> dict:
    """Extract text, notes, and structure from a PowerPoint presentation.

    Returns:
        dict with: slides (list), metadata, total_slides
    """
    try:
        from pptx import Presentation
    except ImportError:
        return {"error": "python-pptx not installed. Run: pip install 'tova[documents]'"}

    prs = Presentation(io.BytesIO(content))

    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        texts.append(row_text)

        # Speaker notes
        notes = ""
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                notes = notes_frame.text.strip()

        # Layout name
        layout = slide.slide_layout.name if slide.slide_layout else ""

        slides.append({
            "slide_number": i + 1,
            "title": texts[0] if texts else "",
            "content": texts[1:] if len(texts) > 1 else [],
            "all_text": "\n".join(texts),
            "notes": notes,
            "layout": layout,
            "shape_count": len(slide.shapes),
        })

    return {
        "format": "pptx",
        "slides": slides,
        "total_slides": len(slides),
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
    }


def create_pptx(
    title: str = "",
    slides: list[dict] | None = None,
) -> bytes:
    """Create a PowerPoint presentation.

    Args:
        title: Presentation title (first slide)
        slides: List of slide definitions:
            [{"title": "Slide Title", "content": ["Bullet 1", "Bullet 2"], "notes": "Speaker notes"}]
            [{"title": "Title", "content": "Plain text body"}]

    Returns:
        Presentation as bytes (.pptx)
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        raise ImportError("python-pptx not installed. Run: pip install 'tova[documents]'")

    prs = Presentation()

    # Title slide
    if title:
        slide_layout = prs.slide_layouts[0]  # Title Slide layout
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = ""

    # Content slides
    for slide_def in (slides or []):
        slide_title = slide_def.get("title", "")
        content = slide_def.get("content", [])
        notes = slide_def.get("notes", "")

        # Use Title + Content layout
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)

        if slide_title and slide.shapes.title:
            slide.shapes.title.text = slide_title

        # Content
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()

            if isinstance(content, str):
                tf.text = content
            elif isinstance(content, list):
                for j, bullet in enumerate(content):
                    if j == 0:
                        tf.text = str(bullet)
                    else:
                        p = tf.add_paragraph()
                        p.text = str(bullet)

        # Speaker notes
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def edit_pptx(content: bytes, edits: list[dict]) -> bytes:
    """Edit an existing PowerPoint presentation.

    Args:
        content: Existing .pptx file bytes
        edits: List of edits:
            [{"action": "edit_slide", "slide": 1, "title": "New Title", "content": ["New bullet"]}]
            [{"action": "add_slide", "title": "New Slide", "content": ["Point 1", "Point 2"]}]
            [{"action": "delete_slide", "slide": 3}]
            [{"action": "edit_notes", "slide": 1, "notes": "Updated speaker notes"}]
            [{"action": "replace_text", "find": "old", "replace": "new"}]

    Returns:
        Modified presentation as bytes
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx not installed. Run: pip install 'tova[documents]'")

    prs = Presentation(io.BytesIO(content))

    for edit in edits:
        action = edit.get("action", "")

        if action == "edit_slide":
            slide_num = edit.get("slide", 1) - 1
            if 0 <= slide_num < len(prs.slides):
                slide = prs.slides[slide_num]
                new_title = edit.get("title")
                new_content = edit.get("content")

                if new_title and slide.shapes.title:
                    slide.shapes.title.text = new_title

                if new_content and len(slide.placeholders) > 1:
                    body = slide.placeholders[1]
                    tf = body.text_frame
                    tf.clear()
                    if isinstance(new_content, str):
                        tf.text = new_content
                    elif isinstance(new_content, list):
                        for j, bullet in enumerate(new_content):
                            if j == 0:
                                tf.text = str(bullet)
                            else:
                                p = tf.add_paragraph()
                                p.text = str(bullet)

        elif action == "add_slide":
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = edit.get("title", "")
            content = edit.get("content", [])

            if title and slide.shapes.title:
                slide.shapes.title.text = title
            if content and len(slide.placeholders) > 1:
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.clear()
                if isinstance(content, list):
                    for j, bullet in enumerate(content):
                        if j == 0:
                            tf.text = str(bullet)
                        else:
                            p = tf.add_paragraph()
                            p.text = str(bullet)

        elif action == "delete_slide":
            slide_num = edit.get("slide", 1) - 1
            if 0 <= slide_num < len(prs.slides):
                rId = prs.slides._sldIdLst[slide_num].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[slide_num]

        elif action == "edit_notes":
            slide_num = edit.get("slide", 1) - 1
            if 0 <= slide_num < len(prs.slides):
                slide = prs.slides[slide_num]
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = edit.get("notes", "")

        elif action == "replace_text":
            find_text = edit.get("find", "")
            replace_text = edit.get("replace", "")
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if find_text in run.text:
                                    run.text = run.text.replace(find_text, replace_text)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── PDF (.pdf) ─────────────────────────────────────────────────────


def read_pdf(content: bytes) -> dict:
    """Extract text from a PDF document.

    Returns:
        dict with: text, pages (list), metadata, total_pages
    """
    try:
        from pypdf import PdfReader
    except ImportError:
        return {"error": "pypdf not installed. Run: pip install 'tova[documents]'"}

    reader = PdfReader(io.BytesIO(content))

    pages = []
    full_text = []

    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        pages.append({
            "page_number": i + 1,
            "text": text,
            "char_count": len(text),
        })
        full_text.append(text)

    metadata = {}
    if reader.metadata:
        metadata = {
            "title": reader.metadata.get("/Title", ""),
            "author": reader.metadata.get("/Author", ""),
            "creator": reader.metadata.get("/Creator", ""),
            "producer": reader.metadata.get("/Producer", ""),
            "created": str(reader.metadata.get("/CreationDate", "")),
        }

    return {
        "format": "pdf",
        "text": "\n\n".join(full_text),
        "pages": pages,
        "total_pages": len(pages),
        "metadata": metadata,
    }


def create_pdf(
    title: str = "",
    content: str = "",
    sections: list[dict] | None = None,
) -> bytes:
    """Create a PDF document.

    Args:
        title: Document title
        content: Plain text content
        sections: Structured sections:
            [{"heading": "Section 1", "text": "Body text..."}]

    Returns:
        PDF as bytes
    """
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.units import inch
    except ImportError:
        raise ImportError("reportlab not installed. Run: pip install 'tova[documents]'")

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4)
    styles = getSampleStyleSheet()
    story = []

    if title:
        story.append(Paragraph(title, styles["Title"]))
        story.append(Spacer(1, 0.3 * inch))

    if content:
        for para in content.split("\n\n"):
            para = para.strip()
            if para:
                story.append(Paragraph(para, styles["Normal"]))
                story.append(Spacer(1, 0.15 * inch))

    if sections:
        for section in sections:
            heading = section.get("heading", "")
            text = section.get("text", "")
            if heading:
                story.append(Paragraph(heading, styles["Heading2"]))
                story.append(Spacer(1, 0.1 * inch))
            if text:
                for para in text.split("\n\n"):
                    para = para.strip()
                    if para:
                        story.append(Paragraph(para, styles["Normal"]))
                        story.append(Spacer(1, 0.1 * inch))

    if not story:
        story.append(Paragraph("Empty document", styles["Normal"]))

    doc.build(story)
    return buf.getvalue()


# ── Format Detection & Routing ─────────────────────────────────────


def detect_format(filename: str) -> str:
    """Detect document format from filename."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    return {
        "docx": "docx", "doc": "docx",
        "xlsx": "xlsx", "xls": "xlsx",
        "pptx": "pptx", "ppt": "pptx",
        "pdf": "pdf",
        "csv": "csv",
        "txt": "text", "md": "text", "json": "json",
    }.get(ext, "unknown")


def read_document(content: bytes, filename: str) -> dict:
    """Read any supported document format.

    Routes to the correct reader based on file extension.
    """
    fmt = detect_format(filename)
    readers = {
        "docx": read_docx,
        "xlsx": read_xlsx,
        "pptx": read_pptx,
        "pdf": read_pdf,
    }

    reader = readers.get(fmt)
    if reader:
        return reader(content)

    # Fallback: try text decode
    if fmt in ("text", "csv", "json"):
        try:
            text = content.decode("utf-8")
            return {"format": fmt, "text": text, "size": len(content)}
        except UnicodeDecodeError:
            return {"format": "binary", "size": len(content), "error": "Cannot decode as text"}

    return {"format": "unknown", "size": len(content), "error": f"Unsupported format: {filename}"}
